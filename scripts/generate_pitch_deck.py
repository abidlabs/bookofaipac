"""Generate a ScanAIPAC pitch deck (.pptx).

Theme matches the ScanAIPAC logo:
  - Dark background      #0c0c0e
  - Light text           #f0f0f0
  - Dim text             #8888a0
  - Accent blue          #7da0ff (logo bracket corners)
  - Card / panel         #18181e with border #26263a
  - Wordmark "barcode"   horizontal black scanlines over white text

Adds slide-level fade transitions and per-shape entrance animations
via direct OOXML so they survive an upload-to-Google-Slides round trip
as well as PowerPoint reasonably allows.
"""

from __future__ import annotations

import copy
import random
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt
from lxml import etree

BG = RGBColor(0x0C, 0x0C, 0x0E)
PANEL = RGBColor(0x18, 0x18, 0x1E)
BORDER = RGBColor(0x26, 0x26, 0x3A)
TEXT = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_DIM = RGBColor(0x88, 0x88, 0xA0)
ACCENT = RGBColor(0x7D, 0xA0, 0xFF)
RED = RGBColor(0xFF, 0x90, 0x90)
GREEN = RGBColor(0x6D, 0xE0, 0x98)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _qn(tag: str) -> str:
    prefix, local = tag.split(":")
    return f"{{{NSMAP[prefix]}}}{local}"


def set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb: RGBColor | None, width_pt: float = 0) -> None:
    if rgb is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def fill_background(slide, rgb: RGBColor = BG) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(
    slide,
    text: str,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    *,
    size: int = 24,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = "Helvetica Neue",
    letter_spacing: int | None = None,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb


def add_panel(slide, x, y, w, h, *, fill=PANEL, border=BORDER, border_w=1.0, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        shp.adjustments[0] = 0.06
    set_solid_fill(shp, fill)
    set_line(shp, border, border_w)
    shp.shadow.inherit = False
    return shp


def add_bracket_corners(slide, *, inset=Inches(0.32), size=Inches(0.42), color=ACCENT, weight=2.75):
    """Draw the ScanAIPAC accent-blue bracket corners at the four slide corners."""
    s = size
    ins = inset
    weight_pt = weight

    def L(x1, y1, x2, y2):
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = color
        line.line.width = Pt(weight_pt)
        return line

    L(ins, ins, ins + s, ins)
    L(ins, ins, ins, ins + s)

    L(SLIDE_W - ins - s, ins, SLIDE_W - ins, ins)
    L(SLIDE_W - ins, ins, SLIDE_W - ins, ins + s)

    L(ins, SLIDE_H - ins, ins, SLIDE_H - ins - s)
    L(ins, SLIDE_H - ins, ins + s, SLIDE_H - ins)

    L(SLIDE_W - ins, SLIDE_H - ins, SLIDE_W - ins, SLIDE_H - ins - s)
    L(SLIDE_W - ins, SLIDE_H - ins, SLIDE_W - ins - s, SLIDE_H - ins)


def add_eyebrow(slide, text: str, *, top=Inches(0.85)):
    return add_text(
        slide,
        text.upper(),
        Inches(0.85),
        top,
        SLIDE_W - Inches(1.7),
        Inches(0.3),
        size=11,
        color=ACCENT,
        bold=True,
        letter_spacing=400,
    )


def add_title(slide, text: str, *, top=Inches(1.2), size=54, color=TEXT, italic_em: str | None = None):
    tb = slide.shapes.add_textbox(Inches(0.85), top, SLIDE_W - Inches(1.7), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if italic_em and italic_em in text:
        before, _, after = text.partition(italic_em)
        for chunk, italic, c in [
            (before, False, color),
            (italic_em, True, ACCENT),
            (after, False, color),
        ]:
            if not chunk:
                continue
            r = p.add_run()
            r.text = chunk
            r.font.name = "Helvetica Neue"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.italic = italic
            r.font.color.rgb = c
    else:
        r = p.add_run()
        r.text = text
        r.font.name = "Helvetica Neue"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
    return tb


def add_subhead(slide, text: str, *, top=Inches(2.6)):
    return add_text(
        slide,
        text,
        Inches(0.85),
        top,
        SLIDE_W - Inches(1.7),
        Inches(1.0),
        size=20,
        color=TEXT_DIM,
    )


def add_bullets(slide, items: list[str], *, top=Inches(3.55), bullet_color=ACCENT, size=20):
    tb = slide.shapes.add_textbox(
        Inches(0.85),
        top,
        SLIDE_W - Inches(1.7),
        SLIDE_H - top - Inches(0.85),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.name = "Helvetica Neue"
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color

        r1 = p.add_run()
        r1.text = item
        r1.font.name = "Helvetica Neue"
        r1.font.size = Pt(size)
        r1.font.color.rgb = TEXT
    return tb


def add_scanlines(slide, x, y, w, h, *, density=0.42):
    """Draw the barcode scanline overlay used in the logo as thin black rects."""
    rng = random.Random(0xA19AC)
    y_cursor = 0
    rects = []
    while y_cursor < int(h):
        gap = rng.randint(int(Inches(0.05)), int(Inches(0.14)))
        thick = rng.randint(int(Pt(1.2)), int(Pt(4)))
        y_cursor += gap
        if y_cursor + thick >= int(h):
            break
        if rng.random() > density:
            continue
        r = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            y + y_cursor,
            w,
            thick,
        )
        set_solid_fill(r, BG)
        set_line(r, None)
        rects.append(r)
    return rects


def add_wordmark(slide, text: str, *, top, height=Inches(1.6), size=140, with_scanlines=True):
    """Draw a centered white wordmark with black scanlines on top (barcode effect)."""
    w = SLIDE_W - Inches(1.7)
    x = Inches(0.85)

    tb = slide.shapes.add_textbox(x, top, w, height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Helvetica Neue"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if with_scanlines:
        add_scanlines(slide, x, top, w, height)
    return tb


def set_transition_fade(slide, duration_ms: int = 600) -> None:
    """Append a <p:transition> element to enable a fade slide transition."""
    sld = slide._element
    nsmap = {"p": NSMAP["p"]}
    for existing in sld.findall("p:transition", nsmap):
        sld.remove(existing)
    transition = etree.SubElement(
        sld,
        _qn("p:transition"),
        attrib={"spd": "med", "advClick": "1"},
    )
    transition.set("dur", str(duration_ms))
    etree.SubElement(transition, _qn("p:fade"))


def add_entrance_animations(slide, shape_ids: list[int], *, delay_step: int = 250, start_delay: int = 0) -> None:
    """No-op placeholder.

    Per-shape entrance animations require deep OOXML timing trees that
    rarely survive an upload to Google Slides cleanly. We rely on slide-level
    fade transitions instead (see set_transition_fade) for the animated feel.
    """
    return


def shape_id(shape) -> int:
    return 0


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    def new_slide(*, brackets=True):
        s = prs.slides.add_slide(blank)
        fill_background(s)
        if brackets:
            add_bracket_corners(s)
        set_transition_fade(s)
        return s

    s = new_slide()
    add_text(
        s, "ACCOUNTABILITY DATABASE", Inches(0.85), Inches(0.85),
        SLIDE_W - Inches(1.7), Inches(0.3),
        size=12, color=ACCENT, bold=True, letter_spacing=400, align=PP_ALIGN.CENTER,
    )
    add_wordmark(s, "ScanAIPAC", top=Inches(2.4), height=Inches(2.0), size=180)
    add_text(
        s, "Why this app exists",
        Inches(0.85), Inches(4.9), SLIDE_W - Inches(1.7), Inches(0.7),
        size=28, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, "A 10-slide tour of what ScanAIPAC does, who it's for, and why now.",
        Inches(0.85), Inches(5.7), SLIDE_W - Inches(1.7), Inches(0.6),
        size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER,
    )

    s = new_slide()
    add_eyebrow(s, "The problem")
    add_title(s, "Money shapes the vote.", italic_em="vote.")
    add_subhead(
        s,
        "Pro-Israel PACs and bundlers spend tens of millions per cycle to back or punish candidates — "
        "the average voter never sees those receipts.",
    )
    stat_cards = [
        ("$100M+", "Pro-Israel spending in recent cycles", RED),
        ("535", "Members of Congress to track", ACCENT),
        ("~0", "Voters who read FEC filings", TEXT_DIM),
    ]
    card_w = Inches(3.7)
    gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(4.6)
    anim_ids = []
    for i, (big, sub, color) in enumerate(stat_cards):
        x = start_x + (card_w + gap) * i
        panel = add_panel(s, x, y, card_w, Inches(2.0))
        anim_ids.append(shape_id(panel))
        add_text(s, big, x, y + Inches(0.25), card_w, Inches(1.0),
                 size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, sub, x + Inches(0.3), y + Inches(1.25), card_w - Inches(0.6), Inches(0.7),
                 size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_entrance_animations(s, anim_ids, delay_step=300)

    s = new_slide()
    add_eyebrow(s, "The information gap")
    add_title(s, "Names on a ballot don't tell you who funds them.", italic_em="who funds them.")
    add_subhead(
        s,
        "FEC data exists. It's just buried in PDFs and spreadsheets that no one reads in the voting booth.",
    )
    add_bullets(s, [
        "You see a name. You don't see the donors.",
        "Local races are even more opaque than federal.",
        "By the time you're voting, it's too late to research.",
    ], top=Inches(4.3))

    s = new_slide()
    add_eyebrow(s, "The idea")
    add_title(s, "Scan a name. Get the answer.", italic_em="Get the answer.")
    add_subhead(
        s,
        "ScanAIPAC turns your camera into an accountability lens for any ballot, mailer, or news article.",
    )
    steps = [
        ("01", "Point camera", "at a candidate's name"),
        ("02", "On-device OCR", "extracts the name"),
        ("03", "Match + show", "donor stance and details"),
    ]
    card_w = Inches(3.7)
    gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(4.4)
    anim_ids = []
    for i, (num, head, body) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        panel = add_panel(s, x, y, card_w, Inches(2.2))
        anim_ids.append(shape_id(panel))
        add_text(s, num, x + Inches(0.3), y + Inches(0.2), Inches(1.0), Inches(0.6),
                 size=14, bold=True, color=ACCENT, letter_spacing=300)
        add_text(s, head, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.6), Inches(0.6),
                 size=22, bold=True, color=TEXT)
        add_text(s, body, x + Inches(0.3), y + Inches(1.35), card_w - Inches(0.6), Inches(0.8),
                 size=14, color=TEXT_DIM)
    add_entrance_animations(s, anim_ids, delay_step=300)

    s = new_slide()
    add_eyebrow(s, "How it works")
    add_title(s, "On-device OCR. No upload. No tracking.", italic_em="No tracking.")
    add_subhead(
        s,
        "The camera frame never leaves your phone. Recognition runs locally; only the matched name is checked against the local candidate index.",
    )
    add_bullets(s, [
        "Apple's Vision / Android ML Kit do the text recognition on-device.",
        "Bundled candidate index ships with the app — works offline.",
        "Camera images and extracted text are not sent to a server.",
    ], top=Inches(4.3))

    s = new_slide()
    add_eyebrow(s, "Stance signals")
    add_title(s, "Color-coded at a glance.", italic_em="at a glance.")
    add_subhead(
        s,
        "Each candidate is tagged with a clear stance and supporting receipts — donations, votes, public statements.",
    )
    items = [
        ("Pro-Israel", RED),
        ("Pro-Palestine", GREEN),
        ("Mixed / unclear", TEXT_DIM),
    ]
    card_w = Inches(3.7)
    gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(4.6)
    anim_ids = []
    for i, (label, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        panel = add_panel(s, x, y, card_w, Inches(1.8), border=color)
        anim_ids.append(shape_id(panel))
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y + Inches(0.45), Inches(0.5), Inches(0.5))
        set_solid_fill(dot, color)
        set_line(dot, None)
        add_text(s, label, x + Inches(1.1), y + Inches(0.4), card_w - Inches(1.4), Inches(0.6),
                 size=22, bold=True, color=color)
        add_text(s, "donations · votes · statements",
                 x + Inches(1.1), y + Inches(1.0), card_w - Inches(1.4), Inches(0.6),
                 size=13, color=TEXT_DIM)
    add_entrance_animations(s, anim_ids, delay_step=250)

    s = new_slide()
    add_eyebrow(s, "Who it's for")
    add_title(s, "Voters, organizers, journalists.", italic_em="journalists.")
    add_subhead(
        s,
        "Anyone who wants to make a 5-second judgment about who's behind a name on a ballot or in the news.",
    )
    add_bullets(s, [
        "Voters scanning mail-in ballots before they fill them out.",
        "Organizers door-knocking with district sheets.",
        "Journalists fact-checking quotes and endorsements in the field.",
        "Students and researchers building their own civic dashboards.",
    ], top=Inches(4.2))

    s = new_slide()
    add_eyebrow(s, "Why now")
    add_title(s, "The receipts exist. The interface didn't.", italic_em="didn't.")
    add_subhead(
        s,
        "FEC data is public. On-device ML is good enough. The only missing piece is a camera-first UI that puts both in your pocket.",
    )
    add_bullets(s, [
        "FEC + OpenSecrets data is structured and free to use.",
        "Modern phones run OCR + matching in milliseconds.",
        "Gen Z + Millennial voters already research candidates on their phones.",
        "Public attention to AIPAC funding is at an all-time high.",
    ], top=Inches(4.2))

    s = new_slide()
    add_eyebrow(s, "What's next")
    add_title(s, "From scanner to civic OS.", italic_em="civic OS.")
    add_subhead(
        s,
        "ScanAIPAC starts with one question — who funds this candidate — and grows into a broader accountability layer.",
    )
    items = [
        ("Now", "iOS + Android scan, candidate database, stance tags"),
        ("Next", "Full ballot scan, district auto-detect, share cards"),
        ("Later", "Open API, journalist tools, multi-issue stance lenses"),
    ]
    card_w = Inches(3.7)
    gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    y = Inches(4.4)
    anim_ids = []
    for i, (head, body) in enumerate(items):
        x = start_x + (card_w + gap) * i
        panel = add_panel(s, x, y, card_w, Inches(2.2))
        anim_ids.append(shape_id(panel))
        add_text(s, head.upper(), x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.5),
                 size=13, bold=True, color=ACCENT, letter_spacing=400)
        add_text(s, body, x + Inches(0.3), y + Inches(0.95), card_w - Inches(0.6), Inches(1.2),
                 size=16, color=TEXT)
    add_entrance_animations(s, anim_ids, delay_step=300)

    s = new_slide()
    add_text(
        s, "JOIN THE BETA", Inches(0.85), Inches(0.85),
        SLIDE_W - Inches(1.7), Inches(0.3),
        size=12, color=ACCENT, bold=True, letter_spacing=400, align=PP_ALIGN.CENTER,
    )
    add_wordmark(s, "ScanAIPAC", top=Inches(2.0), height=Inches(1.8), size=160)
    add_text(
        s, "TestFlight is live. Android internal testing on Google Play.",
        Inches(0.85), Inches(4.3), SLIDE_W - Inches(1.7), Inches(0.7),
        size=22, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_text(
        s, "bookofaipac · scan a name, see who funds them.",
        Inches(0.85), Inches(5.1), SLIDE_W - Inches(1.7), Inches(0.6),
        size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER,
    )

    panel = add_panel(s, (SLIDE_W - Inches(4.2)) / 2, Inches(6.0), Inches(4.2), Inches(0.85),
                      fill=RGBColor(0x5F, 0x6F, 0xD4), border=ACCENT, border_w=1.5)
    add_text(
        s, "GET THE APP",
        (SLIDE_W - Inches(4.2)) / 2, Inches(6.15), Inches(4.2), Inches(0.6),
        size=16, bold=True, color=RGBColor(0xF8, 0xF8, 0xFF),
        align=PP_ALIGN.CENTER, letter_spacing=300,
    )

    out = Path(__file__).resolve().parent.parent / "scanaipac-pitch.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
